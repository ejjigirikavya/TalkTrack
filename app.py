from flask import Flask, render_template, request, redirect
import sqlite3
from pptx import Presentation
import difflib

app = Flask(__name__)

# ---------- DATABASE ----------
def init_db():
    conn = sqlite3.connect("users.db")
    cur = conn.cursor()

    cur.execute("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE,
        password TEXT
    )
    """)

    conn.commit()
    conn.close()


# ---------- PPT TEXT ----------
def extract_ppt_text(file):
    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "

    return text


# ---------- ANALYSIS ----------
filler_words = ["um", "uh", "like", "you know", "basically", "and", "that"]

def calculate_accuracy(ppt_text, spoken_text):
    matcher = difflib.SequenceMatcher(None, ppt_text.lower(), spoken_text.lower())
    return round(matcher.ratio() * 100, 2)


def count_fillers(text):
    words = text.lower().split()
    return sum(word in filler_words for word in words)


# 🔥 FINAL SPEED FIX
def calculate_wpm(text):
    words = len(text.split())

    wpm = words * 3

    if wpm < 60:
        return 60
    elif wpm > 150:
        return 150
    else:
        return wpm


def analyze_speech(text):
    fillers = count_fillers(text)
    wpm = calculate_wpm(text)
    pauses = text.count("...")
    return fillers, wpm, pauses


def generate_ai_feedback(acc, fillers, wpm, pauses):
    feedback = []

    if acc < 50:
        feedback.append("You need to follow PPT more closely.")
    elif acc < 75:
        feedback.append("Good attempt, improve alignment.")
    else:
        feedback.append("Excellent presentation!")

    if fillers > 5:
        feedback.append("Reduce filler words.")
    else:
        feedback.append("Good clarity.")

    if wpm < 100:
        feedback.append("Speak faster.")
    elif wpm > 150:
        feedback.append("Slow down.")
    else:
        feedback.append("Good speed.")

    if pauses > 5:
        feedback.append("Too many pauses.")
    else:
        feedback.append("Good flow.")

    return feedback


# ---------- ROUTES ----------
@app.route('/')
def home():
    return render_template('index.html')


@app.route('/signup')
def signup():
    return render_template('signup.html')


@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')


@app.route('/upload', methods=['POST'])
def upload():
    file = request.files['ppt']
    file.save("uploaded.pptx")
    return redirect('/dashboard')


@app.route('/analyze', methods=['POST'])
def analyze():
    spoken_text = request.form['text']

    if len(spoken_text.strip()) == 0:
        return render_template('dashboard.html', error="No speech detected!", done=False)

    ppt_text = extract_ppt_text("uploaded.pptx")

    accuracy = calculate_accuracy(ppt_text, spoken_text)
    fillers, wpm, pauses = analyze_speech(spoken_text)
    feedback = generate_ai_feedback(accuracy, fillers, wpm, pauses)

    return render_template(
        'dashboard.html',
        spoken_text=spoken_text,
        accuracy=accuracy,
        fillers=fillers,
        wpm=wpm,
        pauses=pauses,
        feedback=feedback,
        done=True
    )


if __name__ == '__main__':
    init_db()
    app.run(debug=True)